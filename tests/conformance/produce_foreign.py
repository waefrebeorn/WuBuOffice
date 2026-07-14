#!/usr/bin/env python3
"""Produce OOXML files with INDEPENDENT libraries (openpyxl, python-pptx) so the
WuBuOffice C readers can be tested against files we did not write ourselves.

Usage: produce_foreign.py <out_dir>
Exit codes:
  0  wrote every file it could; prints one line per file: "<kind> <path>"
  2  a required library is missing -> caller should SKIP (control of destiny is
     still ours; we just can't exercise the foreign path in this environment)
"""
import sys, os

def main():
    if len(sys.argv) < 2:
        print("usage: produce_foreign.py <out_dir>"); return 2
    out = sys.argv[1]
    os.makedirs(out, exist_ok=True)

    try:
        from openpyxl import Workbook
    except Exception:
        print("openpyxl unavailable"); return 2
    try:
        from pptx import Presentation
    except Exception:
        print("python-pptx unavailable"); return 2

    # --- xlsx via openpyxl: multi-sheet, inline formula with NO cached value ---
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws["A1"] = "Item"; ws["B1"] = "Cost"
    ws["A2"] = "Engine"; ws["B2"] = 1200.5
    ws["A3"] = "Docs"; ws["B3"] = 320
    ws["B4"] = "=SUM(B2:B3)"     # openpyxl writes <f> but blanks <v>
    ws2 = wb.create_sheet("Summary")
    ws2["A1"] = "Total"; ws2["B1"] = 1520.5
    xlsx = os.path.join(out, "foreign.xlsx")
    wb.save(xlsx)
    print(f"xlsx {xlsx}")

    # --- pptx via python-pptx: real layouts, placeholder-typed title shapes ---
    prs = Presentation()
    layout = prs.slide_layouts[1]   # Title and Content
    s = prs.slides.add_slide(layout)
    s.shapes.title.text = "Foreign Slide One"
    s.placeholders[1].text = "Alpha bullet.\nBeta bullet."
    s2 = prs.slides.add_slide(layout)
    s2.shapes.title.text = "Foreign Slide Two"
    s2.placeholders[1].text = "Only one."
    pptx = os.path.join(out, "foreign.pptx")
    prs.save(pptx)
    print(f"pptx {pptx}")
    return 0

if __name__ == "__main__":
    sys.exit(main())
